#!/usr/bin/env python3
"""
create_plots_AVE_simul_1_haplo.py
AUTHOR: Dr. Royal Truman
VERSION: 0.9
"""
import sys
import csv
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from collections import defaultdict
from io import BytesIO
import time
from PIL import Image
from pathlib import Path
from typing import Dict, List, Tuple, Any, Optional, Callable, Union

start_time = time.time()

# === Load YAML config file ===
CONFIG_FILE = Path("config_1_haplo_v1.yaml")

try:
    import yaml
except ImportError:
    yaml = None
    print("❌ Required package 'pyyaml' not found.")
    print("👉 Install it with: pip install pyyaml")
    sys.exit(1)

# Check if config file exists
if not CONFIG_FILE.exists():
    print(f"❌ Configuration file '{CONFIG_FILE}' not found in your directory.")
    print(f" Copy it over, and modify the parameters if wished, then rerun.")
    sys.exit(1)

# Load config
try:
    with open(CONFIG_FILE, 'r') as f:
        config = yaml.safe_load(f)
except Exception as e:
    print(f"❌ Error reading or parsing '{CONFIG_FILE}': {e}")
    print("Please check that the file is valid YAML (correct indentation, colons, etc.).")
    sys.exit(1)

# Extract runtime parameters
try:
    CFG_OUTPUT_TYPE = config['OUTPUT_TYPE']
    if CFG_OUTPUT_TYPE not in ['powerpoint', 'images', 'both']:
        raise ValueError("OUTPUT_TYPE must be one of: 'powerpoint', 'images', 'both'")
except KeyError as e:
    print(f"❌ Missing required setting 'OUTPUT_TYPE' in config: {e}")
    print(f"👉 Please make sure '{CONFIG_FILE}' includes all required keys.")
    sys.exit(1)
except ValueError as e:
    print(f"❌ Invalid value in config: {e}")
    sys.exit(1)

CFG_LINE_WIDTH = config.get('LINE_WIDTH', 2.0)
try:
    CFG_LINE_WIDTH = float(CFG_LINE_WIDTH)
except (ValueError, TypeError):
    print(f"⚠️  Invalid LINE_WIDTH '{CFG_LINE_WIDTH}' in config. Using default 2.0.")
    CFG_LINE_WIDTH = 2.0

CFG_SIM_COLORS = config.get('SIM_COLORS', {})
if not isinstance(CFG_SIM_COLORS, dict):
    print("⚠️  SIM_COLORS in config is not a dictionary. Ignoring and using defaults.")
    CFG_SIM_COLORS = {}

# =============================
# CONSTANTS
# =============================

# --- Input Files ---
INPUT_FILE_A = Path("haplo_AVE_A_per_gen_1H.txt")
INPUT_FILE_a = Path("haplo_AVE_a__per_gen_1H.txt")
INPUT_FILE_SLIDE2 = Path("A_haplo_AVE_vals_per_gen_1H.txt")
PARAMS_FILE = Path("in_1_haplo.txt")

# --- Output Files ---
PPTX_OUTPUT = Path("AVE_per_gen_1_haplo.pptx")
TIFF_DPI = 600

# --- Slide 1 Output Files ---
SLIDE1_FREQ_A_FILE = Path("AVE_1_haplo_freq_A_new.tiff")
SLIDE1_N_A_FILE = Path("AVE_1_haplo_N_A_new.tiff")
SLIDE1_FREQ_a_FILE = Path("AVE_1_haplo_freq_a_existing.tiff")
SLIDE1_N_a_FILE = Path("AVE_1_haplo_N_a_existing.tiff")
SLIDE1_COMBO_FILE = Path("AVE_1_haplo_combo_1.tiff")

# --- Slide 2 Output Files ---
SLIDE2_FREQ_FILE = Path("AVE_1_haplo_freq.tiff")
SLIDE2_HETEROZ_FILE = Path("AVE_1_haplo_heteroz.tiff")
SLIDE2_HOMOZYG_FILE = Path("AVE_1_haplo_homozyg.tiff")
SLIDE2_POPUL_FILE = Path("AVE_1_haplo_popul_size.tiff")
SLIDE2_COMBO_FILE = Path("AVE_1_haplo_combo_2.tiff")

# --- Slide Titles ---
SLIDE1_TITLE_TEXT = "Average frequencies of 'A', and 'a' haplotypes, and average population size vs effective generations"
SLIDE2_TITLE_TEXT = "Frequencies of 'A', and 'a' haplotypes, heterozygosity and homozygosity over generations"

# --- Title Slide Image ---
TITLE_IMAGE_PATH = Path("AVE_PER_GEN_1_haplo_IMAGE.PNG")

# ========================
# Utility Functions (SHARED)
# ========================

def validate_file(filepath: Path) -> None:
    """Validate that a file exists and is readable"""
    if not filepath.exists():
        raise FileNotFoundError(f"Input file '{filepath}' not found in current directory.")
    if not filepath.is_file():
        raise IsADirectoryError(f"'{filepath}' is a directory, not a file.")


def load_data(filepath: Path) -> Tuple[Dict[int, List[Dict[str, Any]]], List[str]]:
    """
    Load the CSV data, parse it, and group by SimNr.
    Returns a dict: {sim_nr: list of rows as dicts}
    """
    data: Dict[int, List[Dict[str, Any]]] = defaultdict(list)
    with open(filepath, 'r', newline='', encoding='utf-8') as f:
        reader = csv.DictReader(f, delimiter=';')
        for row in reader:
            for key, val in row.items():
                try:
                    row[key] = float(val)
                except (ValueError, TypeError):
                    pass
            sim_nr = int(row['SimNr'])
            data[sim_nr].append(row)
    return data, reader.fieldnames


def extract_series(data_dict: Dict[int, List[Dict[str, Any]]], x_key: str, y_key: str) -> Dict[int, Tuple[List[float], List[float]]]:
    """
    Extracts x and y series grouped by SimNr.
    Returns: {sim_nr: (x_list, y_list)}
    """
    series: Dict[int, Tuple[List[float], List[float]]] = {}
    for sim_nr, rows in data_dict.items():
        x_vals = [row[x_key] for row in rows]
        y_vals = [row[y_key] for row in rows]
        series[sim_nr] = (x_vals, y_vals)
    return series


def format_with_commas(x: float, pos: int) -> str:
    """Formatter for tick labels to include comma separators."""
    return f"{int(x):,}"


def get_sim_colors(all_sim_nrs: List[int]) -> Dict[int, Any]:
    """Get colors for simulations from config with fallback"""
    colors: Dict[int, Any] = {}
    cmap = plt.get_cmap('tab10')  # fallback palette

    for sim in all_sim_nrs:
        if sim in CFG_SIM_COLORS:
            colors[sim] = CFG_SIM_COLORS[sim]
        else:
            print(f"⚠️  Warning: No color defined in config for SimNr {sim}. Using fallback.")
            colors[sim] = cmap((sim - 1) % 10)  # deterministic fallback

    return colors


def _plot_series_on_ax(
    ax: plt.Axes,
    series_dict: Dict[int, Tuple[List[float], List[float]]],
    sim_colors: Dict[int, Any],
    linestyle: str = '-',
    label_suffix: str = '',
    linewidth: float = CFG_LINE_WIDTH
) -> None:
    """Helper to plot simulation series on given axis — avoids duplication."""
    for sim_nr in sorted(series_dict.keys()):
        x_vals, y_vals = series_dict[sim_nr]
        label = f"Sim {sim_nr}{label_suffix}"
        ax.plot(x_vals, y_vals,
                color=sim_colors[sim_nr],
                linewidth=linewidth,
                linestyle=linestyle,
                label=label)


def get_max_generation(series_list: List[Dict[int, Tuple[List[float], List[float]]]]) -> float:
    """Centralized helper to compute max generation across multiple series."""
    all_gens: List[float] = []
    for series in series_list:
        for x_vals, _ in series.values():
            all_gens.extend(x_vals)
    return max(all_gens) if all_gens else 0.0


# ========================
# REUSABLE HELPER FUNCTIONS
# ========================

def create_single_plot(
    series_dict: Dict[int, Tuple[List[float], List[float]]],
    sim_colors: Dict[int, Any],
    max_gen: float,
    xlabel: str,
    ylabel: str,
    title: Optional[str] = None,
    legend_ncol: int = 1,
    extra_series_list: Optional[List[Dict[str, Any]]] = None
) -> plt.Figure:
    """
    Creates a single matplotlib figure plotting multiple simulation series.
    """
    fig, ax = plt.subplots(figsize=(8, 6), dpi=TIFF_DPI)
    
    # Plot main series
    _plot_series_on_ax(ax, series_dict, sim_colors, '-', '', CFG_LINE_WIDTH)
    
    # Plot extra series (if any)
    if extra_series_list:
        for extra in extra_series_list:
            extra_dict = extra['series_dict']
            linestyle = extra.get('linestyle', '--')
            label_suffix = extra.get('label_suffix', '')
            _plot_series_on_ax(ax, extra_dict, sim_colors, linestyle, label_suffix, CFG_LINE_WIDTH)

    ax.set_xlim(0, max_gen)
    ax.set_xlabel(xlabel, fontsize=14)
    ax.set_ylabel(ylabel, fontsize=14)
    ax.tick_params(labelsize=14)
    ax.legend(fontsize=12, loc='best', ncol=legend_ncol)
    
    if max_gen > 999:
        ax.xaxis.set_major_formatter(plt.FuncFormatter(format_with_commas))
    
    if title:
        ax.set_title(title, fontsize=16)
    
    return fig


def create_combo_plot(
    plot_configs: List[Dict[str, Any]],
    sim_colors: Dict[int, Any],
    max_gen: float
) -> plt.Figure:
    """
    Creates a 2x2 combo plot.
    """
    fig, axes = plt.subplots(2, 2, figsize=(16, 12), dpi=TIFF_DPI)
    axes = axes.flatten()

    for idx, config in enumerate(plot_configs):
        ax = axes[idx]
        series_dict = config['series']
        xlabel = config['xlabel']
        ylabel = config['ylabel']
        legend_ncol = config.get('legend_ncol', 1)
        extra_series_list = config.get('extra_series_list', None)

        # Plot main series
        _plot_series_on_ax(ax, series_dict, sim_colors, '-', '', CFG_LINE_WIDTH)

        # Plot extra series
        if extra_series_list:
            for extra in extra_series_list:
                extra_dict = extra['series_dict']
                linestyle = extra.get('linestyle', '--')
                label_suffix = extra.get('label_suffix', '')
                _plot_series_on_ax(ax, extra_dict, sim_colors, linestyle, label_suffix, CFG_LINE_WIDTH)

        ax.set_xlim(0, max_gen)
        ax.set_xlabel(xlabel, fontsize=14)
        ax.set_ylabel(ylabel, fontsize=14)
        ax.tick_params(labelsize=14)
        ax.legend(fontsize=12, loc='best', ncol=legend_ncol)

        if max_gen > 999:
            ax.xaxis.set_major_formatter(plt.FuncFormatter(format_with_commas))

    plt.tight_layout()
    return fig


def save_figures_as_tiffs(figures: List[plt.Figure], filenames: List[Path]) -> None:
    """Save a list of matplotlib figures as TIFF files."""
    for fig, fname in zip(figures, filenames):
        fig.savefig(fname, dpi=TIFF_DPI, bbox_inches='tight')
        plt.close(fig)
    print("Files saved:", [str(f) for f in filenames])


def add_combo_slide_to_powerpoint(prs: Presentation, combo_fig: plt.Figure, title_text: str) -> None:
    """
    Adds a slide with a title and a centered combo plot image.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    left = Inches(1)
    top = Inches(0.2)
    width = Inches(8)
    height = Inches(0.225)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = False
    title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Inches(0.2)
    p.alignment = PP_ALIGN.CENTER

    # Add combo graph
    img_top = Inches(0.7)
    img_left = Inches(0.5)
    img_width = Inches(9)
    img_height = Inches(6.5)

    img_buffer = BytesIO()
    combo_fig.savefig(img_buffer, format='png', dpi=TIFF_DPI, bbox_inches='tight')
    img_buffer.seek(0)

    slide.shapes.add_picture(img_buffer, left=img_left, top=img_top,
                             width=img_width, height=img_height)
    img_buffer.close()
    plt.close(combo_fig)


def load_and_scale_image(
    img_path: Path,
    max_width_in: float,
    max_height_in: float
) -> Optional[Tuple[BytesIO, float, float]]:
    """
    Load image, scale to fit within max dimensions while preserving aspect ratio.
    Returns (buffer, scaled_width, scaled_height) or None if failed.
    """
    try:
        with Image.open(img_path) as img:
            img_width_px, img_height_px = img.size
            img_aspect_ratio = img_width_px / img_height_px

            # Scale image to fit within available space
            if img_aspect_ratio > (max_width_in / max_height_in):
                new_width = max_width_in
                new_height = max_width_in / img_aspect_ratio
            else:
                new_height = max_height_in
                new_width = max_height_in * img_aspect_ratio

            # Save to buffer
            img_buffer = BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            return img_buffer, new_width, new_height

    except Exception as e:
        print(f"⚠️  Warning: Could not load or process '{img_path}': {e}. Skipping image on title slide.")
        return None


# ========================
# TITLE SLIDE: Slide 0
# ========================
def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    # === Top centered title textbox ===
    top_title_text = "1-haplotype simulations. Generated with the Sim-generations suite"
    top_title_left = (slide_width - 7) / 2
    top_title_top = 0.5
    top_title_width = 7
    top_title_height = 0.3

    title_box = slide.shapes.add_textbox(Inches(top_title_left), Inches(top_title_top), Inches(top_title_width), Inches(top_title_height))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    
    run = p.add_run()
    run.text = "1-haplotype simulations. Generated with the "
    run.font.size = Pt(18)
    run.font.italic = False

    run = p.add_run()
    run.text = "Sim-generations"
    run.font.size = Pt(18)
    run.font.italic = True

    run = p.add_run()
    run.text = " suite"
    run.font.size = Pt(18)
    run.font.italic = False

    p.alignment = PP_ALIGN.CENTER

    # === Bottom-left creator info ===
    creator_text = "Created with create_plots_AVE_simul_1_haplo.py (v. 0.9)"
    creator_left = 0.5
    creator_bottom_margin = 0.2
    creator_top = slide_height - creator_bottom_margin - 0.3
    creator_width = 4.5
    creator_height = 0.3

    creator_box = slide.shapes.add_textbox(Inches(creator_left), Inches(creator_top), Inches(creator_width), Inches(creator_height))
    creator_frame = creator_box.text_frame
    creator_frame.clear()
    p = creator_frame.paragraphs[0]
    p.text = creator_text
    p.font.size = Inches(0.15)
    p.font.bold = False

    # === Bottom-right date and time ===
    timestamp = time.strftime("%b. %d, %Y  %H:%M")
    timestamp_right_margin = 1
    timestamp_left = slide_width - timestamp_right_margin - 1
    timestamp_top = creator_top
    timestamp_width = 1
    timestamp_height = 0.3

    timestamp_box = slide.shapes.add_textbox(Inches(timestamp_left), Inches(timestamp_top), Inches(timestamp_width), Inches(timestamp_height))
    timestamp_frame = timestamp_box.text_frame
    timestamp_frame.clear()
    p = timestamp_frame.paragraphs[0]
    p.text = timestamp
    p.font.size = Inches(0.15)
    p.font.bold = False

    # === Centered image ===
    top_margin = top_title_top + top_title_height + 1
    bottom_margin = creator_top - 1
    left_margin = 1
    right_margin = slide_width - 1

    avail_width = right_margin - left_margin
    avail_height = bottom_margin - top_margin

    result = load_and_scale_image(TITLE_IMAGE_PATH, avail_width, avail_height)
    if result:
        img_buffer, new_width, new_height = result
        img_left = left_margin + (avail_width - new_width) / 2
        img_top = top_margin + (avail_height - new_height) / 2
        slide.shapes.add_picture(img_buffer, Inches(img_left), Inches(img_top),
                                 Inches(new_width), Inches(new_height))
        img_buffer.close()


# ========================
# SLIDE 1: Average frequencies of A and a haplotypes
# ========================
def load_slide1_data() -> Dict[str, Dict[int, Tuple[List[float], List[float]]]]:
    """
    Loads data from two CSV files for slide 1.
    Returns dict with keys: 'freq_A', 'N_A', 'freq_a', 'N_a'
    """
    files = {
        'A': INPUT_FILE_A,
        'a': INPUT_FILE_a
    }

    data_dict = {}
    for key, filepath in files.items():
        validate_file(filepath)
        data, _ = load_data(filepath)
        data_dict[key] = data

    series = {}
    series['freq_A'] = extract_series(data_dict['A'], 'EffGen', 'freq_A')
    series['N_A'] = extract_series(data_dict['A'], 'EffGen', 'N')
    series['freq_a'] = extract_series(data_dict['a'], 'EffGen', 'freq_a')
    series['N_a'] = extract_series(data_dict['a'], 'EffGen', 'N')

    return series


def create_slide1_plots(series: Dict[str, Dict[int, Tuple[List[float], List[float]]]], max_gen: float) -> List[plt.Figure]:
    all_sim_nrs = sorted(series['freq_A'].keys())
    colors = get_sim_colors(all_sim_nrs)

    figures = []

    plot_specs = [
        {'series': series['freq_A'], 'xlabel': 'Effective Generations', 'ylabel': "Freq. of 'A' (new haplotype)", 'legend_ncol': 1},
        {'series': series['N_A'], 'xlabel': 'Effective Generations', 'ylabel': 'Population Size (N) - new haplotype', 'legend_ncol': 1},
        {'series': series['freq_a'], 'xlabel': 'Effective Generations', 'ylabel': "Freq. of 'a' (existing haplotype)", 'legend_ncol': 1},
        {'series': series['N_a'], 'xlabel': 'Effective Generations', 'ylabel': 'Population Size (N) - existing haplotype', 'legend_ncol': 1}
    ]

    for spec in plot_specs:
        fig = create_single_plot(
            series_dict=spec['series'],
            sim_colors=colors,
            max_gen=max_gen,
            xlabel=spec['xlabel'],
            ylabel=spec['ylabel'],
            legend_ncol=spec.get('legend_ncol', 1)
        )
        figures.append(fig)

    combo_config = [
        {'series': series['freq_A'], 'xlabel': 'Effective Generations', 'ylabel': "Freq. of 'A'", 'legend_ncol': 2},
        {'series': series['N_A'], 'xlabel': 'Effective Generations', 'ylabel': 'Population Size (N)', 'legend_ncol': 1},
        {'series': series['freq_a'], 'xlabel': 'Effective Generations', 'ylabel': "Freq. of 'a'", 'legend_ncol': 2},
        {'series': series['N_a'], 'xlabel': 'Effective Generations', 'ylabel': 'Population Size (N)', 'legend_ncol': 1}
    ]

    combo_fig = create_combo_plot(combo_config, colors, max_gen)
    figures.append(combo_fig)

    return figures


# ========================
# SLIDE 2: Frequencies, heterozygosity and homozygosity
# ========================
def load_slide2_data() -> Dict[int, List[Dict[str, Any]]]:
    validate_file(INPUT_FILE_SLIDE2)
    data, _ = load_data(INPUT_FILE_SLIDE2)
    return data


def create_slide2_plots(data: Dict[int, List[Dict[str, Any]]], max_gen: float) -> List[plt.Figure]:
    freq_A_series = extract_series(data, 'EffGen', 'freq A')
    freq_a_series = extract_series(data, 'EffGen', 'freq a')
    hetero_series = extract_series(data, 'EffGen', 'hetero')
    homoz_series = extract_series(data, 'EffGen', 'homoz')
    n_series = extract_series(data, 'EffGen', 'N')

    all_sim_nrs = sorted(data.keys())
    colors = get_sim_colors(all_sim_nrs)

    figures = []

    fig1 = create_single_plot(
        series_dict=freq_A_series,
        sim_colors=colors,
        max_gen=max_gen,
        xlabel='Effective Generations',
        ylabel='Freq. A and a haplotypes',
        legend_ncol=2,
        extra_series_list=[{
            'series_dict': freq_a_series,
            'linestyle': '--',
            'label_suffix': ' (a)'
        }]
    )
    figures.append(fig1)

    for series, ylabel in [
        (hetero_series, 'Freq. heterozygous in Aa'),
        (homoz_series, 'Freq. homozygous in Aa'),
        (n_series, 'Population size (N)')
    ]:
        fig = create_single_plot(
            series_dict=series,
            sim_colors=colors,
            max_gen=max_gen,
            xlabel='Effective Generations',
            ylabel=ylabel,
            legend_ncol=1
        )
        figures.append(fig)

    combo_config = [
        {
            'series': freq_A_series,
            'xlabel': 'Effective Generations',
            'ylabel': 'Freq. A and a haplotypes',
            'legend_ncol': 2,
            'extra_series_list': [{
                'series_dict': freq_a_series,
                'linestyle': '--',
                'label_suffix': ' (a)'
            }]
        },
        {'series': n_series, 'xlabel': 'Effective Generations', 'ylabel': 'Population size (N)', 'legend_ncol': 1},
        {'series': hetero_series, 'xlabel': 'Effective Generations', 'ylabel': 'Freq. heterozygous in Aa', 'legend_ncol': 1},
        {'series': homoz_series, 'xlabel': 'Effective Generations', 'ylabel': 'Freq. homozygous in Aa', 'legend_ncol': 1}
    ]

    combo_fig = create_combo_plot(combo_config, colors, max_gen)
    figures.append(combo_fig)

    return figures


# ========================
# SLIDE 3: Parameters and simulation information
# =======================
def add_slide3_parameters(prs: Presentation) -> None:
    validate_file(PARAMS_FILE)
    try:
        with open(PARAMS_FILE, 'r', newline='', encoding='utf-8') as f:
            reader = csv.reader(f, delimiter=';')
            headers = next(reader)
            raw_rows = list(reader)
    except Exception as e:
        raise RuntimeError(f"Failed to read '{PARAMS_FILE}': {e}")

    headers_with_sim = ['Sim Nr'] + headers
    data_with_sim = [[str(idx + 1)] + [cell.strip() for cell in row] for idx, row in enumerate(raw_rows)]

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.225))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = False
    title_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    p = title_frame.paragraphs[0]
    p.text = "Simulation runs and parameters used"
    p.font.size = Inches(0.2)
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.0), Inches(8), Inches(1))
    desc_frame = desc_box.text_frame
    desc_frame.clear()
    desc_frame.word_wrap = True
    desc_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    description_text = """Meaning of the parameters
Sim Nr: Simulation number. Each simulation uses the same parameter values. 
Ni: Initial size of the population at the beginning of the simulation. 1 and 1,000,000,000.
r: Growth rate of the population per generation. Usually ≥ 0, but values as long as > -1.
K: Carrying capacity, the maximum population size in an environment. Must be ≥ Ni.
s_A: Selectivity coefficient for haplotype A. Between -2 and #2.
h_A: Dominance coefficient, i.e., fitness of the heterozygous genotype Aa vs. AA and aa. (Fitness AA = 1 + h × s_A).
p_A_i: Initial proportion of haplotype A in the population. p_a_i is by definition 1 - p_A_i.  Between 0 and 1 
attempts: The number of times each simulation is to be repeat. Each is individual rerun Rep (repetitions) random times."""

    lines = description_text.strip().split('\n')
    first_para = desc_frame.paragraphs[0]
    first_para.text = lines[0]
    first_para.font.size = Inches(0.16)

    for line in lines[1:]:
        p = desc_frame.add_paragraph()
        if ':' in line:
            var_part, rest = line.split(':', 1)
            run = p.add_run()
            run.text = var_part + ':'
            run.font.bold = True
            run = p.add_run()
            run.text = rest
        else:
            p.text = line
        p.font.size = Inches(0.16)

    para_count = len(lines)
    font_in = 0.16
    line_height_in = font_in * 1.2
    estimated_desc_height = para_count * line_height_in + 0.25
    desc_box.height = Inches(estimated_desc_height)

    # Dynamic table
    def estimate_text_width_inches(text: str, font_size_inches: float, is_bold: bool = False) -> float:
        if not text:
            return 0.3
        char_width_ratio = 0.55 if not is_bold else 0.62
        estimated_width = len(str(text)) * font_size_inches * char_width_ratio
        padding = 0.20
        return estimated_width + padding

    num_cols = len(headers_with_sim)
    num_rows = len(data_with_sim) + 1
    header_font_size = 0.195
    data_font_size = 0.16
    col_widths_inches = []

    for col_idx in range(num_cols):
        header_text = headers_with_sim[col_idx]
        header_width = estimate_text_width_inches(header_text, header_font_size, is_bold=True)
        max_data_width = max(
            (estimate_text_width_inches(str(row[col_idx]), data_font_size) if col_idx < len(row) else 0)
            for row in data_with_sim
        )
        required_width = max(header_width, max_data_width, 0.5)
        col_widths_inches.append(required_width)

    actual_table_width_inches = sum(col_widths_inches)
    row_height_inches = 0.30
    actual_table_height_inches = num_rows * row_height_inches

    slide_width_inches = 10
    table_left = Inches((slide_width_inches - actual_table_width_inches) / 2)
    table_top = desc_box.top + desc_box.height + Inches(0.3)

    table = slide.shapes.add_table(
        num_rows, 
        num_cols, 
        table_left, 
        table_top, 
        Inches(actual_table_width_inches), 
        Inches(actual_table_height_inches)
    ).table

    for col_idx in range(num_cols):
        table.columns[col_idx].width = int(Inches(col_widths_inches[col_idx]))

    for row_idx in range(num_rows):
        table.rows[row_idx].height = int(Inches(row_height_inches))

    for j, header in enumerate(headers_with_sim):
        cell = table.cell(0, j)
        cell.text = header
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Inches(0.195)
        para.alignment = PP_ALIGN.CENTER

    for i, row in enumerate(data_with_sim):
        for j, val in enumerate(row):
            if j < len(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Inches(0.16)
                para.alignment = PP_ALIGN.CENTER


# ========================
# Main Execution
# ========================
def main() -> None:
    # === Validate all input files upfront ===
    REQUIRED_INPUT_FILES = [
        INPUT_FILE_A,
        INPUT_FILE_a,
        INPUT_FILE_SLIDE2,
        PARAMS_FILE,
    ]
    for f in REQUIRED_INPUT_FILES:
        validate_file(f)

    # ========================
    # SLIDE 1
    # ========================
    print("📊 Processing Slide 1: Average frequencies of A and a haplotypes...")

    slide1_series = load_slide1_data()

    sim_nrs = set(slide1_series['freq_A'].keys())
    if not all(sim_nrs == set(slide1_series[k].keys()) for k in ['freq_a', 'N_A', 'N_a']):
        raise ValueError("Mismatch in simulation numbers across slide 1 input files.")

    max_gen_slide1 = get_max_generation([
        slide1_series['freq_A'],
        slide1_series['N_A'],
        slide1_series['freq_a'],
        slide1_series['N_a']
    ])

    slide1_figures = create_slide1_plots(slide1_series, max_gen_slide1)
    slide1_combo_fig = slide1_figures[-1]

    if CFG_OUTPUT_TYPE in ['images', 'both']:
        save_figures_as_tiffs(slide1_figures, [
            SLIDE1_FREQ_A_FILE,
            SLIDE1_N_A_FILE,
            SLIDE1_FREQ_a_FILE,
            SLIDE1_N_a_FILE,
            SLIDE1_COMBO_FILE
        ])

    # ========================
    # SLIDE 2
    # ========================
    print("📊 Processing Slide 2: Frequencies, heterozygosity and homozygosity...")

    slide2_data = load_slide2_data()
    if not slide2_data:
        raise ValueError("No data loaded from slide 2 input file.")

    max_gen_slide2 = get_max_generation([
        extract_series(slide2_data, 'EffGen', 'freq A'),
        extract_series(slide2_data, 'EffGen', 'freq a'),
        extract_series(slide2_data, 'EffGen', 'hetero'),
        extract_series(slide2_data, 'EffGen', 'homoz'),
        extract_series(slide2_data, 'EffGen', 'N')
    ])

    slide2_figures = create_slide2_plots(slide2_data, max_gen_slide2)
    slide2_combo_fig = slide2_figures[-1]

    if CFG_OUTPUT_TYPE in ['images', 'both']:
        save_figures_as_tiffs(slide2_figures, [
            SLIDE2_FREQ_FILE,
            SLIDE2_HETEROZ_FILE,
            SLIDE2_HOMOZYG_FILE,
            SLIDE2_POPUL_FILE,
            SLIDE2_COMBO_FILE
        ])

    # ========================
    # Create PowerPoint
    # ========================
    if CFG_OUTPUT_TYPE in ['powerpoint', 'both']:
        prs = Presentation()
        add_title_slide(prs)
        add_combo_slide_to_powerpoint(prs, slide1_combo_fig, SLIDE1_TITLE_TEXT)
        add_combo_slide_to_powerpoint(prs, slide2_combo_fig, SLIDE2_TITLE_TEXT)
        add_slide3_parameters(prs)
        prs.save(PPTX_OUTPUT)
        print(f"💾 PowerPoint saved as '{PPTX_OUTPUT}'.")


# ========================
# Run the program
# ========================
if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"\n❌ Error: {e}")
        sys.exit(1)
    finally:
        current_dir = Path.cwd()
        print(f"📁 Output files stored in: {current_dir}")

        end_time = time.time()
        execution_time = end_time - start_time
        print(f"⏱️ Total runtime: {execution_time:.2f} seconds")
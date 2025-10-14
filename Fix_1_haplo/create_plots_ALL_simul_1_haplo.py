#!/usr/bin/env python3
"""
create_plots_ALL_simul_1_haplo.py
AUTHOR: Dr. Royal Truman
VERSION: 0.9
"""

import sys
import time
from datetime import datetime
from pathlib import Path
from typing import List, Optional

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
from PIL import Image

# ==============================
# CONFIGURATION & CONSTANTS
# ==============================

# Files
CONFIG_FILE = Path("config_1_haplo_v1.yaml")
INPUT_DATA_FILE = Path("out_1_haplo_per_gen.txt")
TITLE_IMAGE_PATH = Path("ALL_PER_GEN_1_haplo_IMAGE.PNG")
POWERPOINT_OUTPUT = Path("ALL_per_gen_1_haplo.pptx")

# Plotting Variables
Y_VARS = ['N', 'freq_A_and_a', 'freq_het_hom']
Y_LABELS = [
    'Population Size (N)',
    "Freq. A and a haplotypes",
    "Freq. heteroz., homoz."
]

# Column names for input data
COLUMN_NAMES = [
    'SimNr', 'Rep', 'attempt', 'Ni', 'r', 'K', 's_A', 'h_A', 'p_A_i',
    'attempts', 'generation', 'N', 'freq_A', 'freq_Aa', 'freq_a', 'homoz'
]

# Numeric columns for conversion
NUMERIC_COLS = [
    'SimNr', 'Rep', 'attempt', 'Ni', 'r', 'K', 's_A', 'h_A',
    'p_A_i', 'attempts', 'generation', 'N', 'freq_A', 'freq_Aa',
    'freq_a', 'homoz'
]

# Slide layout & dimensions (in inches)
TITLE_LAYOUT = {
    'title_top': Inches(0.5),
    'title_left': Inches(1.5),
    'title_width': Inches(9),
    'title_height': Inches(0.3),
    'title_font_size_pt': 14,

    'bottom_left_top_offset': Inches(0.5),
    'bottom_right_top_offset': Inches(0.5),
    'bottom_left_left': Inches(1),
    'bottom_right_right_margin': Inches(1),
    'bottom_box_width': Inches(5),
    'bottom_box_height': Inches(0.3),
    'bottom_font_size_pt': 12,

    'image_top_margin': Inches(1.8),
    'image_bottom_margin': Inches(1.5),
    'image_side_margin': Inches(1),
}

SLIDE_LAYOUT = {
    'title_top': Inches(0.1),
    'title_left': Inches(2),
    'title_width': Inches(4.5),
    'title_height': Inches(0.3),
    'title_font_size_in': Inches(0.16),

    'param_top': Inches(0.3),
    'param_left': Inches(1.0),
    'param_width': Inches(8.0),
    'param_height': Inches(0.4),
    'param_font_size_in': Inches(0.16),

    'plot_left': Inches(0.1),
    'plot_top': Inches(0.8),
    'plot_width_offset': Inches(0.2),
    'plot_height_offset': Inches(1.0),
}

# Plot styling
PLOT_COLORS = {
    'A': 'darkblue',
    'a': 'red',
    'het': 'darkgreen',
    'hom': 'orange',
    'default': 'black'
}

# ==============================
# HELPER FUNCTIONS
# ==============================

def log_error(message: str):
    print(f"❌ {message}")

def log_warning(message: str):
    print(f"⚠️  {message}")

def log_info(message: str):
    print(f"📄 {message}")

def format_large_number(n: float) -> str:
    """Format large numbers with commas (e.g., 1000000 -> '1,000,000')"""
    return f"{int(n):,}"

def add_parameter_box(ax, row0):
    """Add a parameter info box to the center of the plot."""
    Ni_display = int(row0['Ni']) if pd.notna(row0['Ni']) and float(row0['Ni']).is_integer() else row0['Ni']
    K_display = int(row0['K']) if pd.notna(row0['K']) and float(row0['K']).is_integer() else row0['K']
    param_text = (
        f"Initial pop. (Ni) = {Ni_display}\n"
        f"Pop. growth rate (r) = {row0['r']:.3f}\n"
        f"Maximum pop. (K) = {K_display}\n"
        f"Selectivity coefficient s(A) = {row0['s_A']:.3f}\n"
        f"Dominance coefficient, h(A) = {row0['h_A']:.3f}\n"
        f"Proportion A initially = {row0['p_A_i']:.3f}"
    )
    xlim = ax.get_xlim()
    ylim = ax.get_ylim()
    x_mid = (xlim[0] + xlim[1]) / 2
    y_mid = (ylim[0] + ylim[1]) / 2
    ax.text(x_mid, y_mid, param_text,
            fontsize=10,
            ha='left', va='center',
            bbox=dict(boxstyle="round,pad=0.6", facecolor="wheat", alpha=0.9),
            transform=ax.transData)

def plot_trajectories(ax, group: pd.DataFrame, y_var: str, ylabel: str, show_legend: bool = True, show_param_box: bool = False, row0: Optional[pd.Series] = None):
    """
    Plot trajectories for a given y-variable on the given axis.
    Handles special cases for freq_A_and_a and freq_het_hom.
    """
    group['trajectory_id'] = group['Rep'].astype(str) + '_' + group['attempt'].astype(str)

    if y_var == 'freq_A_and_a':
        for _, traj_data in group.groupby('trajectory_id'):
            traj_sorted = traj_data.sort_values('generation')
            ax.plot(traj_sorted['generation'], traj_sorted['freq_A'], color=PLOT_COLORS['A'], linewidth=1.2, label="haplotype 'A'")
            ax.plot(traj_sorted['generation'], traj_sorted['freq_a'], color=PLOT_COLORS['a'], linewidth=1.2, label="haplotype 'a'")
        if show_legend:
            handles, labels = ax.get_legend_handles_labels()
            by_label = dict(zip(labels, handles))
            ax.legend(by_label.values(), by_label.keys(), fontsize=12, loc='best')
    elif y_var == 'freq_het_hom':
        for _, traj_data in group.groupby('trajectory_id'):
            traj_sorted = traj_data.sort_values('generation')
            ax.plot(traj_sorted['generation'], traj_sorted['freq_Aa'], color=PLOT_COLORS['het'], linewidth=1.2, label='heterozygous')
            ax.plot(traj_sorted['generation'], traj_sorted['homoz'], color=PLOT_COLORS['hom'], linewidth=1.2, label='homozygous')
        if show_legend:
            handles, labels = ax.get_legend_handles_labels()
            by_label = dict(zip(labels, handles))
            ax.legend(by_label.values(), by_label.keys(), fontsize=12, loc='best')
    else:
        for _, traj_data in group.groupby('trajectory_id'):
            traj_sorted = traj_data.sort_values('generation')
            ax.plot(traj_sorted['generation'], traj_sorted[y_var], color=PLOT_COLORS['default'], linewidth=1.2)

    ax.set_ylabel(ylabel, fontsize=12)
    ax.grid(True, alpha=0.3)
    ax.tick_params(axis='both', which='major', labelsize=10)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: format_large_number(x)))

    if show_param_box and row0 is not None:
        add_parameter_box(ax, row0)

# ==============================
# POWERPOINT FUNCTIONS
# ==============================

def add_title_slide(presentation: Presentation):
    """Add title slide as the first slide in the presentation."""
    blank_slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(blank_slide_layout)

    # Move this slide to the beginning
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(0, slides[-1])

    # Top title textbox
    title_textbox = slide.shapes.add_textbox(
        TITLE_LAYOUT['title_left'],
        TITLE_LAYOUT['title_top'],
        TITLE_LAYOUT['title_width'],
        TITLE_LAYOUT['title_height']
    )
    title_frame = title_textbox.text_frame
    title_frame.text = "1-haplotype simulations, all trajectories. Generated with the Sim-generations suite"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Inches(TITLE_LAYOUT['title_font_size_pt'] / 72)
    title_para.alignment = PP_ALIGN.CENTER

    # Make "Sim-generations" italic
    runs = title_para.runs
    for run in runs:
        text = run.text
        if "Sim-generations" in text:
            before_sim = text.split("Sim-generations")[0]
            after_sim = text.split("Sim-generations")[1] if len(text.split("Sim-generations")) > 1 else ""

            run.text = before_sim

            italic_run = title_para.add_run()
            italic_run.text = "Sim-generations"
            italic_run.font.italic = True
            italic_run.font.size = Inches(TITLE_LAYOUT['title_font_size_pt'] / 72)

            if after_sim:
                normal_run = title_para.add_run()
                normal_run.text = after_sim
                normal_run.font.size = Inches(TITLE_LAYOUT['title_font_size_pt'] / 72)

    # Bottom left textbox
    bottom_left_textbox = slide.shapes.add_textbox(
        TITLE_LAYOUT['bottom_left_left'],
        presentation.slide_height - TITLE_LAYOUT['bottom_left_top_offset'],
        TITLE_LAYOUT['bottom_box_width'],
        TITLE_LAYOUT['bottom_box_height']
    )
    bottom_left_frame = bottom_left_textbox.text_frame
    bottom_left_frame.text = "Created with create_plots_ALL_simul_1_haplo.py (v. 0.9)"
    bottom_left_para = bottom_left_frame.paragraphs[0]
    bottom_left_para.font.size = Inches(TITLE_LAYOUT['bottom_font_size_pt'] / 72)
    bottom_left_para.alignment = PP_ALIGN.LEFT

    # Bottom right textbox
    current_time = datetime.now()
    date_time_str = current_time.strftime("%b. %d, %Y  %H:%M")

    bottom_right_textbox = slide.shapes.add_textbox(
        presentation.slide_width - TITLE_LAYOUT['bottom_right_right_margin'] - TITLE_LAYOUT['bottom_box_width'],
        presentation.slide_height - TITLE_LAYOUT['bottom_right_top_offset'],
        TITLE_LAYOUT['bottom_box_width'],
        TITLE_LAYOUT['bottom_box_height']
    )
    bottom_right_frame = bottom_right_textbox.text_frame
    bottom_right_frame.text = date_time_str
    bottom_right_para = bottom_right_frame.paragraphs[0]
    bottom_right_para.font.size = Inches(TITLE_LAYOUT['bottom_font_size_pt'] / 72)
    bottom_right_para.alignment = PP_ALIGN.RIGHT

    # Add center image with aspect ratio preservation
    if TITLE_IMAGE_PATH.exists():
        try:
            with Image.open(TITLE_IMAGE_PATH) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height

            available_top = TITLE_LAYOUT['image_top_margin']
            available_bottom = TITLE_LAYOUT['image_bottom_margin']
            available_height = presentation.slide_height - available_top - available_bottom

            available_left = TITLE_LAYOUT['image_side_margin']
            available_right = TITLE_LAYOUT['image_side_margin']
            available_width = presentation.slide_width - available_left - available_right

            if available_width / available_height > aspect_ratio:
                img_display_height = available_height
                img_display_width = available_height * aspect_ratio
            else:
                img_display_width = available_width
                img_display_height = available_width / aspect_ratio

            img_left = (presentation.slide_width - img_display_width) / 2
            img_top = available_top + (available_height - img_display_height) / 2

            slide.shapes.add_picture(str(TITLE_IMAGE_PATH), img_left, img_top, img_display_width, img_display_height)

        except Exception as e:
            log_warning(f"Could not add image {TITLE_IMAGE_PATH}: {e}")
    else:
        log_warning(f"Image file {TITLE_IMAGE_PATH} not found. Continuing without image.")

# ==============================
# CORE FUNCTION: Create Slide for Simulation
# ==============================

def create_slide_for_simulation(
    sim_nr: int,
    df: pd.DataFrame,
    presentation: Optional[Presentation],
    tmpdir: str,
    y_vars: List[str],
    y_labels: List[str],
    output_type: str
):
    """
    Create plots and PowerPoint slide (if applicable) for a single simulation number.
    Also saves TIFF images if output_type includes 'images'.
    """
    # Filter and sort data for this simulation
    group = df[df['SimNr'] == sim_nr].copy()
    group = group.sort_values('generation')

    if group.empty:
        log_warning(f"No data found for simulation {sim_nr}")
        return

    rep_val = int(group['Rep'].max())
    attempts_val = int(group['attempt'].max())
    total_runs = rep_val * attempts_val

    row0 = group.iloc[0]

    # Warn if parameters vary within simulation
    param_cols = ['Ni', 'r', 'K', 's_A', 'h_A', 'p_A_i']
    for col in param_cols:
        if group[col].nunique() > 1:
            log_warning(f"Sim {sim_nr} has varying '{col}' values. Using first row.")

    # Slide Title
    title_text = (
        f"Simulation {sim_nr} "
        f"(Repetitions = {rep_val}, attempts = {attempts_val}, total runs = {total_runs})"
    )

    # Parameter Text (comma-separated, single line)
    Ni_display = int(row0['Ni']) if pd.notna(row0['Ni']) and float(row0['Ni']).is_integer() else row0['Ni']
    K_display = int(row0['K']) if pd.notna(row0['K']) and float(row0['K']).is_integer() else row0['K']
    param_text = (
        f"Initial pop. (Ni) = {Ni_display}, "
        f"pop. growth rate (r) = {row0['r']:.3f}, "
        f"maximum pop. (K) = {K_display}, "
        f"selectivity coefficient s(A) = {row0['s_A']:.3f}, "
        f"dominance coefficient, h(A) = {row0['h_A']:.3f}, "
        f"proportion A initially = {row0['p_A_i']:.3f}."
    )

    # CREATE COMPOSITE FIGURE FOR POWERPOINT
    if output_type in ['powerpoint', 'both']:
        fig, axes = plt.subplots(3, 1, figsize=(12, 8.4), sharex=True)

        for i, (y_var, ylabel) in enumerate(zip(y_vars, y_labels)):
            plot_trajectories(axes[i], group, y_var, ylabel, show_legend=True)

        axes[-1].set_xlabel('generations', fontsize=12)
        plt.tight_layout(rect=(0, 0, 1, 0.97))

        img_path = Path(tmpdir) / f"sim_{sim_nr}.png"
        fig.savefig(img_path, dpi=600, bbox_inches='tight')
        plt.close(fig)

        # ADD TO POWERPOINT
        blank_slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_slide_layout)

        # Slide Title
        title_box = slide.shapes.add_textbox(
            SLIDE_LAYOUT['title_left'],
            SLIDE_LAYOUT['title_top'],
            SLIDE_LAYOUT['title_width'],
            SLIDE_LAYOUT['title_height']
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        p = title_frame.paragraphs[0]
        p.font.size = SLIDE_LAYOUT['title_font_size_in']
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # Parameter Textbox
        text_box = slide.shapes.add_textbox(
            SLIDE_LAYOUT['param_left'],
            SLIDE_LAYOUT['param_top'],
            SLIDE_LAYOUT['param_width'],
            SLIDE_LAYOUT['param_height']
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = param_text
        p = text_frame.paragraphs[0]
        p.font.size = SLIDE_LAYOUT['param_font_size_in']
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.LEFT
        p.space_before = 0
        p.space_after = 0
        p.line_spacing = 1.0

        # Add Image
        img_left = SLIDE_LAYOUT['plot_left']
        img_top = SLIDE_LAYOUT['plot_top']
        img_width = presentation.slide_width - SLIDE_LAYOUT['plot_width_offset']
        img_height = presentation.slide_height - SLIDE_LAYOUT['plot_height_offset']

        slide.shapes.add_picture(str(img_path), img_left, img_top, img_width, img_height)

    # SAVE .tiff IMAGES
    if output_type in ['images', 'both']:
        for y_var, ylabel in zip(y_vars, y_labels):
            fig_single, ax = plt.subplots(figsize=(10, 6))
            ax.set_xlabel('generations', fontsize=12)
            ax.set_title(f"Simulation {sim_nr} - {ylabel}", fontsize=14)

            plot_trajectories(ax, group, y_var, ylabel, show_legend=(y_var in ['freq_A_and_a', 'freq_het_hom']), show_param_box=True, row0=row0)

            # Generate filename
            if y_var == 'N':
                img_filename = f"ALL_1_haplo_popul_size_sim_{sim_nr}.tiff"
            elif y_var == 'freq_A_and_a':
                img_filename = f"ALL_1_haplo_freq_sim_{sim_nr}.tiff"
            elif y_var == 'freq_het_hom':
                img_filename = f"ALL_1_haplo_het_hom_sim_{sim_nr}.tiff"
            else:
                img_filename = f"Simulation_{sim_nr}_{y_var}.tiff"

            img_path_single = Path(img_filename)

            try:
                fig_single.savefig(img_path_single, format='tiff', dpi=600, bbox_inches='tight')
            except Exception as e:
                log_error(f"Error saving {img_path_single}: {e}")
            finally:
                plt.close(fig_single)

        print(f"💾 Files saved for Sim. {sim_nr}: "
              f"ALL_1_haplo_popul_size_sim_{sim_nr}.tiff "
              f"ALL_1_haplo_freq_sim_{sim_nr}.tiff "
              f"ALL_1_haplo_het_hom_sim_{sim_nr}.tiff")

# ==============================
# MAIN EXECUTION
# ==============================

if __name__ == "__main__":
    start_time = time.time()

    # === Load YAML config file ===
    try:
        import yaml
    except ImportError:
        log_error("Required package 'pyyaml' not found.")
        print("👉 Install it with: pip install pyyaml")
        sys.exit(1)

    if not CONFIG_FILE.exists():
        log_error(f"Configuration file '{CONFIG_FILE}' not found in your directory.")
        print("Copy it over, and modify the parameters if wished, then rerun.")
        sys.exit(1)

    try:
        with open(CONFIG_FILE, 'r') as f:
            config = yaml.safe_load(f)
            log_info(f"Loading configuration from '{CONFIG_FILE}'...")
    except Exception as e:
        log_error(f"Error reading or parsing '{CONFIG_FILE}': {e}")
        print("Please check that the file is valid YAML (correct indentation, colons, etc.).")
        sys.exit(1)

    # Extract and validate OUTPUT_TYPE
    try:
        OUTPUT_TYPE = config['OUTPUT_TYPE']
        if OUTPUT_TYPE not in ['powerpoint', 'images', 'both']:
            raise ValueError("OUTPUT_TYPE must be one of: 'powerpoint', 'images', 'both'")
    except KeyError as e:
        log_error(f"Missing required setting 'OUTPUT_TYPE' in config: {e}")
        print(f"👉 Please make sure '{CONFIG_FILE}' includes all required keys.")
        sys.exit(1)
    except ValueError as e:
        log_error(f"Invalid value in config: {e}")
        sys.exit(1)

    # Ensure Pillow is installed
    try:
        from PIL import Image
    except ImportError:
        raise ImportError(
            "Pillow is required to save TIFF images.\n"
            "Install it using: pip install pillow"
        )

    # === Load and clean data ===
    if not INPUT_DATA_FILE.exists():
        raise FileNotFoundError(f"Input data file not found: {INPUT_DATA_FILE}")

    log_info(f"Reading simulation data from '{INPUT_DATA_FILE}'...")
    df = pd.read_csv(INPUT_DATA_FILE, sep=None, engine='python', names=COLUMN_NAMES, header=None)

    df[NUMERIC_COLS] = df[NUMERIC_COLS].apply(pd.to_numeric, errors='coerce')
    df.dropna(subset=['SimNr', 'generation'], inplace=True)
    df['SimNr'] = df['SimNr'].astype(int)

    sim_nrs = sorted(df['SimNr'].unique())

    # === Initialize PowerPoint if needed ===
    presentation = None
    if OUTPUT_TYPE in ['powerpoint', 'both']:
        presentation = Presentation()
        add_title_slide(presentation)

    # === Generate outputs ===
    with tempfile.TemporaryDirectory() as tmpdir:
        for sim_nr in sim_nrs:
            create_slide_for_simulation(sim_nr, df, presentation, tmpdir, Y_VARS, Y_LABELS, OUTPUT_TYPE)

        if presentation:
            presentation.save(POWERPOINT_OUTPUT)
            log_info(f"PowerPoint saved as '{POWERPOINT_OUTPUT}'")

    current_dir = Path.cwd()
    print(f"📁 Output files stored in: {current_dir}")

    execution_time = time.time() - start_time
    print(f"⏱️ Total runtime: {execution_time:.2f} seconds")
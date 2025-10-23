#!/usr/bin/env python3
"""
create_plots_AVE_simul_2_haplos.py
AUTHOR: Dr. Royal Truman
VERSION: 0.9
"""
import os
import sys
import time
import tempfile
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    import yaml
except ImportError:
    yaml = None
    print("❌ Required package 'pyyaml' not found.")
    print("👉 Install it with: pip install pyyaml")
    sys.exit(1)


# ======================
# CONSTANTS
# ======================
CONFIG_FILE = "config_2_haplos_v1.yaml"
pptx_filename = 'AVE_per_gen_2_haplos.pptx'
DPI = 600

# Font sizes (can be made configurable later via cfg)
LABEL_FONTSIZE = 12
TICK_FONTSIZE = 12
TITLE_FONTSIZE = 14

# Input filenames
files = {
    'gene_A': 'A_haplo_AVE_vals_per_gen_2H.txt',
    'gene_B': 'B_haplo_AVE_vals_per_gen_2H.txt',
    'pan_heteroz': 'Pan-heteroz_AVE_vals_per_gen_2H.txt',
    'pan_homoz': 'Pan-homoz_AVE_vals_per_gen_2H.txt'
}

# TIFF output mapping
TIFF_NAME_MAP = {
    'graph1': 'AVE_A_haplo_freq.tiff',
    'graph2': 'AVE_A_haplo_heteroz.tiff',
    'graph3': 'AVE_B_haplo_freq.tiff',
    'graph4': 'AVE_B_haplo_heteroz.tiff',
    'graph5': 'AVE_Pan-heteroz.tiff',
    'graph6': 'AVE_Pan-homoz.tiff',
    'graph7': 'AVE_Popul_size.tiff'
}


# ======================
# CONFIGURATION
# ======================
# Load and validate config
try:
    with open(CONFIG_FILE, 'r') as f:
        raw_config = yaml.safe_load(f)
        print(f"📄 Loading configuration from '{CONFIG_FILE}'...")
except FileNotFoundError:
    print(f"❌ Configuration file '{CONFIG_FILE}' not found in your directory.")
    print("👉 Copy it over, and modify the parameters if wished, then rerun.")
    sys.exit(1)
except Exception as e:
    print(f"❌ Error reading or parsing '{CONFIG_FILE}': {e}")
    print("Please check that the file is valid YAML (correct indentation, colons, etc.).")
    sys.exit(1)


class ConfigManager:
    """
    Centralized configuration access and validation.
    Ensures required keys exist and provides typed access with defaults.
    """
    def __init__(self, config_dict):
        self._config = config_dict
        self._validate()

    def _validate(self):
        required_keys = ['OUTPUT_TYPE']
        for key in required_keys:
            if key not in self._config:
                raise KeyError(f"❌ Missing required config key: '{key}'")

        output_type = self._config['OUTPUT_TYPE']
        if output_type not in ['powerpoint', 'images', 'both']:
            raise ValueError(f"❌ Invalid OUTPUT_TYPE: '{output_type}'. Must be one of: 'powerpoint', 'images', 'both'")

    def get_output_type(self):
        return self._config['OUTPUT_TYPE']

    def get_line_width(self, default=2.0):
        return self._config.get('LINE_WIDTH', default)

    def get_sim_colors(self, default=None):
        if default is None:
            default = {}
        return self._config.get('SIM_COLORS', default)


# Initialize config manager
try:
    cfg = ConfigManager(raw_config)
    OUTPUT_TYPE = cfg.get_output_type()
except (KeyError, ValueError) as e:
    print(f"❌ Configuration error: {e}")
    sys.exit(1)


# ======================
# HELPER FUNCTIONS
# ======================

# ----------------------
# Plotting Helpers
# ----------------------
def format_thousands_axis(ax):
    """
    Format the x-axis of the given axes to display numbers with comma separators (e.g., 1,000, 2,500).
    Only applies to non-negative values.
    """
    ax.get_xaxis().set_major_formatter(
        plt.FuncFormatter(lambda x, _: f'{int(x):,}' if x >= 0 else '')
    )


def apply_plot_style(ax, xlabel="", ylabel=""):
    """
    Apply consistent styling to a matplotlib Axes object:
    - Set axis labels with standard font size
    - Set tick label size
    - Apply legend styling (if legend exists)
    - Apply x-axis thousands formatting
    """
    ax.set_xlabel(xlabel, fontsize=LABEL_FONTSIZE)
    ax.set_ylabel(ylabel, fontsize=LABEL_FONTSIZE)
    ax.tick_params(axis='both', which='major', labelsize=TICK_FONTSIZE)

    legend = ax.get_legend()
    if legend is None:
        handles, labels = ax.get_legend_handles_labels()
        if handles:
            ax.legend(handles, labels, loc='best', fontsize=8,
                      frameon=True, fancybox=False, edgecolor='black')
    else:
        for text in legend.get_texts():
            text.set_fontsize(8)# Here the font size for all legends is set
        legend.set_frame_on(True)
        legend.set_fancybox(False)
        legend.set_edgecolor('black')

    format_thousands_axis(ax)


def plot_scatter_lines(ax, df, x_col, y_cols, group_col, title, xlabel, ylabel, cmap='tab10'):
    """
    Plot lines for each group (e.g., Sim Nr), with optional custom colors and line styles.
    """
    sim_nrs = df[group_col].unique()
    try:
        colormap = plt.colormaps.get_cmap(cmap)
    except ValueError:
        colormap = plt.colormaps.get_cmap('tab10')

    sim_colors = cfg.get_sim_colors()
    line_width = cfg.get_line_width()

    for idx, sim in enumerate(sim_nrs):
        group_df = df[df[group_col] == sim]
        color = sim_colors.get(str(sim), sim_colors.get(int(sim), colormap(idx % colormap.N)))
        for y_col in y_cols:
            if y_col in group_df.columns:
                x_vals = group_df[x_col]
                y_vals = group_df[y_col]
                
                dashed_alleles = {'Freq a', 'Freq b'}
                linestyle = 'dashed' if y_col in dashed_alleles else 'solid'

                legend_label = f"Sim {sim} – {y_col}"
                ax.plot(x_vals, y_vals, color=color, linewidth=line_width,
                        linestyle=linestyle, label=legend_label)

    apply_plot_style(ax, xlabel, ylabel)


def create_plots(data):
    """Create individual plots for each gene and pan-statistic."""
    plots = {}

    plot_specs = [
        ('graph1', data['gene_A'], 'Gen', ['Freq A', 'Freq a'], 'Sim Nr', 'Freq. A and a haplotypes'),
        ('graph2', data['gene_A'], 'Gen', ['Freq Aa'], 'Sim Nr', 'Freq. heterozygous in Aa'),
        ('graph3', data['gene_B'], 'Gen', ['Freq B', 'Freq b'], 'Sim Nr', 'Freq. B and b haplotypes'),
        ('graph4', data['gene_B'], 'Gen', ['Freq Bb'], 'Sim Nr', 'Freq. heterozygous in Bb'),
        ('graph5', data['pan_heteroz'], 'Gen', ['Pan heteroz'], 'Sim Nr', 'Freq. pan heterozygous'),
        ('graph6', data['pan_homoz'], 'Gen', ['Pan homoz'], 'Sim Nr', 'Freq. pan homozygous'),
        ('graph7', data['pan_homoz'], 'Gen', ['N'], 'Sim Nr', 'Population Size (N)'),
    ]

    for key, df, x_col, y_cols, group_col, ylabel in plot_specs:
        fig, ax = plt.subplots(figsize=(10, 5))
        plot_scatter_lines(ax, df, x_col, y_cols, group_col,
                           title=None, xlabel='Generations', ylabel=ylabel)
        ax.set_title("")  # Explicitly clear title
        plots[key] = fig

    return plots


# ----------------------
# PowerPoint Helpers
# ----------------------
def add_slide_title(slide, text, left=3, top=0.2, width=8, height=0.3,
                    alignment=PP_ALIGN.LEFT, bold=True, fontsize=TITLE_FONTSIZE,
                    vertical_anchor=None):
    """
    Add or update a styled title on a PowerPoint slide.
    """
    title_box = slide.shapes.title
    if title_box is None:
        title_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

    title_box.left = Inches(left)
    title_box.top = Inches(top)
    title_box.width = Inches(width)
    title_box.height = Inches(height)

    text_frame = title_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    if vertical_anchor is not None:
        text_frame.vertical_anchor = vertical_anchor

    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p.text = text
    p.alignment = alignment

    for run in p.runs:
        run.font.size = Pt(fontsize)
        run.font.bold = bold


def add_plot_to_ppt(presentation, fig, title="", save_as_tiff=None):
    """Add a matplotlib figure to a PowerPoint slide."""
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    add_slide_title(slide, title)

    # Save to temp file
    temp_img = tempfile.NamedTemporaryFile(delete=False, suffix='.tiff')
    fig.savefig(temp_img.name, dpi=DPI, format='tiff', bbox_inches='tight')
    plt.close(fig)

    # Add image to slide
    img_left = Inches(0.5)
    img_top = Inches(0.5)
    img_width = Inches(12)
    slide.shapes.add_picture(temp_img.name, img_left, img_top, width=img_width)
    temp_img.close()

    if save_as_tiff:
        fig.savefig(save_as_tiff, dpi=DPI, bbox_inches='tight')
        print(f"💾 Saved image: {save_as_tiff}")


# ----------------------
# Utility Helpers
# ----------------------
def validate_files(file_dict):
    """Check that all required input files exist."""
    missing = [filename for filename in file_dict.values() if not os.path.exists(filename)]
    if missing:
        raise FileNotFoundError(f"The following files are missing: {', '.join(missing)}")


def load_data(file_dict):
    """Load all data files into a dictionary of DataFrames."""
    data = {}
    for key, filename in file_dict.items():
        df = pd.read_csv(filename, delimiter=';', dtype={'Sim Nr': str})
        df['Gen'] = pd.to_numeric(df['Gen'], errors='coerce')
        data[key] = df
    return data


def save_figure(fig, filename):
    """Save figure to file with standard DPI and layout."""
    fig.savefig(filename, dpi=DPI, bbox_inches='tight')
    plt.close(fig)
    print(f"...{filename}")


# ======================
# SLIDE CREATION FUNCTIONS (IN ORDER)
# ======================

def create_slide_1_title(presentation):
    """Add decorative title slide with logo and metadata."""
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)
    slide_width = presentation.slide_width.inches
    slide_height = presentation.slide_height.inches

    # Top text
    top_text = "2-haplotype simulations, average results. Generated with the Sim-generations suite"
    top_left = (slide_width - 9) / 2
    top_top = 0.5
    top_width = 9
    top_height = 0.3
    top_txBox = slide.shapes.add_textbox(Inches(top_left), Inches(top_top), Inches(top_width), Inches(top_height))
    tf_top = top_txBox.text_frame
    tf_top.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_top.clear()
    p_top = tf_top.paragraphs[0]
    p_top.space_before = Pt(0)
    p_top.space_after = Pt(0)
    p_top.alignment = PP_ALIGN.CENTER
    parts = top_text.split("Sim-generations")
    run1 = p_top.add_run()
    run1.text = parts[0]
    run1.font.size = Pt(14)
    run_italic = p_top.add_run()
    run_italic.text = "Sim-generations"
    run_italic.font.size = Pt(14)
    run_italic.font.italic = True
    run2 = p_top.add_run()
    run2.text = parts[1]
    run2.font.size = Pt(14)

    # Bottom left
    bl_text = "Created with create_plots_AVE_simul_2_haplos.py (v. 0.9)"
    bl_left = 1
    bottom_margin_inches = 0.5 / 2.54
    bl_top = slide_height - bottom_margin_inches - 0.3
    bl_width = 5
    bl_height = 0.3
    bl_txBox = slide.shapes.add_textbox(Inches(bl_left), Inches(bl_top), Inches(bl_width), Inches(bl_height))
    tf_bl = bl_txBox.text_frame
    tf_bl.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_bl.clear()
    p_bl = tf_bl.paragraphs[0]
    p_bl.space_before = Pt(0)
    p_bl.space_after = Pt(0)
    p_bl.alignment = PP_ALIGN.LEFT
    run_bl = p_bl.add_run()
    run_bl.text = bl_text
    run_bl.font.size = Pt(12)

    # Bottom right
    date_str = time.strftime("%b. %d, %Y  %H:%M", time.localtime())
    date_str = date_str.replace("Sep.", "Sept.")
    br_left = slide_width - 1 - 5
    br_top = bl_top
    br_width = 5
    br_height = 0.3
    br_txBox = slide.shapes.add_textbox(Inches(br_left), Inches(br_top), Inches(br_width), Inches(br_height))
    tf_br = br_txBox.text_frame
    tf_br.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf_br.clear()
    p_br = tf_br.paragraphs[0]
    p_br.space_before = Pt(0)
    p_br.space_after = Pt(0)
    p_br.alignment = PP_ALIGN.RIGHT
    run_br = p_br.add_run()
    run_br.text = date_str
    run_br.font.size = Pt(12)

    # Image
    image_path = "AVE_PER_GEN_2_HAPLOS_IMAGE.PNG"
    if not os.path.exists(image_path):
        print(f"Warning: Image file '{image_path}' not found. Continuing without adding the image.")
    else:
        temp_shape = slide.shapes.add_picture(image_path, Inches(0), Inches(0))
        aspect = temp_shape.width / temp_shape.height
        temp_shape._element.getparent().remove(temp_shape._element)

        top_margin_min = 1.8
        bottom_max = bl_top - 1
        max_h_centered = 7.5 - 3.6
        max_w = slide_width - 2
        max_h = min(max_h_centered, bottom_max - top_margin_min)

        img_w = max_w
        img_h = img_w / aspect
        if img_h > max_h:
            img_h = max_h
            img_w = img_h * aspect

        img_left = (slide_width - img_w) / 2
        img_top = (slide_height - img_h) / 2
        slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), Inches(img_w))


def create_slide_2_haplotype_frequencies_and_pop_size(presentation, data):
    """Combine Freq A/a, Freq B/b (side by side) and Population Size (full width below)."""
    fig = plt.figure(figsize=(14, 8.5))
    gs = fig.add_gridspec(2, 2, height_ratios=[1, 0.7], hspace=0.3, wspace=0.3)

    ax1 = fig.add_subplot(gs[0, 0])
    plot_scatter_lines(ax1, data['gene_A'], 'Gen', ['Freq A', 'Freq a'], 'Sim Nr',
                       title=None, xlabel='Generations', ylabel='Freq. A and a haplotypes')

    ax2 = fig.add_subplot(gs[0, 1])
    plot_scatter_lines(ax2, data['gene_B'], 'Gen', ['Freq B', 'Freq b'], 'Sim Nr',
                       title=None, xlabel='Generations', ylabel='Freq. B and b haplotypes')

    ax3 = fig.add_subplot(gs[1, :])
    plot_scatter_lines(ax3, data['pan_homoz'], 'Gen', ['N'], 'Sim Nr',
                       title=None, xlabel='Generations', ylabel='Population Size (N)')

    fig.subplots_adjust(top=0.90, bottom=0.08, left=0.08, right=0.97)

    if OUTPUT_TYPE == 'both':
        save_figure(fig, 'AVE_combo_graph_Nr_1.tiff')

    add_plot_to_ppt(
        presentation,
        fig,
        title="Haplotype Frequencies and Population Size over Generations"
    )


def create_slide_3_heterozygosity_and_pan_stats_grid(presentation, data):
    """Combine heterozygosity and pan-stats in 2x2 grid."""
    fig, ((ax1, ax2), (ax3, ax4)) = plt.subplots(2, 2, figsize=(14, 8),
                                                 gridspec_kw={'hspace': 0.3, 'wspace': 0.3})

    plot_specs = [
        (ax1, data['gene_A'], 'Gen', ['Freq Aa'], 'Sim Nr', 'Freq. heterozygous in Aa'),
        (ax2, data['gene_B'], 'Gen', ['Freq Bb'], 'Sim Nr', 'Freq. heterozygous in Bb'),
        (ax3, data['pan_heteroz'], 'Gen', ['Pan heteroz'], 'Sim Nr', 'Freq. pan heterozygous'),
        (ax4, data['pan_homoz'], 'Gen', ['Pan homoz'], 'Sim Nr', 'Freq. pan homozygous'),
    ]

    for ax, df, x_col, y_cols, group_col, ylabel in plot_specs:
        plot_scatter_lines(ax, df, x_col, y_cols, group_col,
                           title=None, xlabel='Generations', ylabel=ylabel)

    fig.subplots_adjust(top=0.92, bottom=0.1, left=0.08, right=0.97, hspace=0.3, wspace=0.3)

    if OUTPUT_TYPE == 'both':
        save_figure(fig, 'AVE_combo_graph_Nr_2.tiff')

    add_plot_to_ppt(
        presentation,
        fig,
        title="Frequencies of A/a and B/b heterozygosity, pan homozygosity, and pan heterozygosity"
    )


def create_slide_4_parameters_table(presentation):
    """Add slide with simulation parameter definitions and input data table."""
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    add_slide_title(
        slide,
        "Simulation runs and parameters used",
        left=3, top=0.2, width=5, height=0.3,
        alignment=PP_ALIGN.CENTER,
        fontsize=14,
        vertical_anchor=MSO_ANCHOR.MIDDLE
    )

    # --- Text box with parameter meanings ---
    text_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(9), Inches(2.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()

    meanings = [
        "          === Meaning of the simulation parameters ===",
        "Sim Nr: Simulation number. Each simulation uses the same parameter values.",
        "Ni: Initial size of the population at the beginning of the simulation. Between 1 and 1,000,000,000.",
        "r: Growth rate of the population per generation. Usually ≥ 0, but can be negative as long as > -1.",
        "K: Carrying capacity, the maximum population size in that environment. Must be ≥ Ni.",
        "s_A: Selectivity coefficient for haplotype A. Between -2 and #2 inclusive.",
        "s_B: Selectivity coefficient for haplotype B. Between -2 and #2 inclusive.",
        "h_A: Dominance coefficient, i.e., fitness of the heterozygous genotype Aa vs. AA and aa. (Fitness AA = 1 + h × s_A).",
        "h_B: Dominance coefficient, i.e., fitness of the heterozygous genotype Bb vs. BB and bb. (Fitness BB = 1 + h × s_B).",
        "p_A_i: Initial proportion of haplotype A in the population. p_a_i is by definition 1 - p_A_i. Between 0 and 1.",
        "p_B_i: Initial proportion of haplotype B in the population. p_b_i is by definition 1 - p_B_i. Between 0 and 1.",
        "attempts: The number of times each simulation is to be rerun. Each is individual rerun Rep (repetitions) random times.",
        " ",
        "          Parameters currently defined in input file in_2_haplos.txt:"
    ]

    for i, line in enumerate(meanings):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.font.size = Pt(12)

        if line.startswith(("Sim Nr:", "Ni:", "r:", "K:", "s_A:", "s_B:", "h_A:", "h_B:", "p_A_i:", "p_B_i:", "attempts:")):
            parts = line.split(":", 1)
            run1 = p.add_run()
            run1.text = parts[0] + ":"
            run1.font.bold = True
            run1.font.size = Pt(12)
            if len(parts) > 1:
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.size = Pt(12)
        else:
            p.text = line

    # --- Load and display input data table ---
    input_csv = 'in_2_haplos.txt'
    if not os.path.exists(input_csv):
        raise FileNotFoundError(f"Required file '{input_csv}' not found.")

    try:
        df = pd.read_csv(input_csv, delimiter=';', dtype=str)
    except Exception as e:
        raise ValueError(f"Failed to read '{input_csv}': {e}")

    df.insert(0, 'Sim Nr', [str(i+1) for i in range(len(df))])

    # --- Calculate dynamic column widths ---
    def estimate_text_width_inches(text, font_size_pt, is_bold=False):
        if not text:
            return 0.3
        font_size_inches = font_size_pt / 72.0
        char_width_ratio = 0.55 if not is_bold else 0.62
        estimated_width = len(str(text)) * font_size_inches * char_width_ratio
        padding = 0.20
        return estimated_width + padding

    cols = len(df.columns)
    rows = len(df) + 1
    header_font_size = 14
    data_font_size = 12

    col_widths_inches = []
    for col_idx, col_name in enumerate(df.columns):
        header_width = estimate_text_width_inches(col_name, header_font_size, is_bold=True)
        max_data_width = max(
            estimate_text_width_inches(str(df.iloc[row_idx][col_name]), data_font_size)
            for row_idx in range(len(df))
        )
        required_width = max(header_width, max_data_width, 0.5)
        col_widths_inches.append(required_width)

    actual_table_width_inches = sum(col_widths_inches)
    row_height_inches = 0.12
    actual_table_height_inches = rows * row_height_inches

    slide_width_inches = 13.33
    table_left = Inches((slide_width_inches - actual_table_width_inches) / 2)
    table_top = text_box.top + text_box.height + Inches(0.1)

    table = slide.shapes.add_table(
        rows, cols, table_left, table_top,
        Inches(actual_table_width_inches), Inches(actual_table_height_inches)
    ).table

    for col_idx in range(cols):
        table.columns[col_idx].width = int(Inches(col_widths_inches[col_idx]))
    for row_idx in range(rows):
        table.rows[row_idx].height = int(Inches(row_height_inches))

    # Fill header
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.text = col_name
        p.font.size = Pt(header_font_size)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Fill data
    for row_idx in range(len(df)):
        for col_idx, col_name in enumerate(df.columns):
            cell = table.cell(row_idx + 1, col_idx)
            tf = cell.text_frame
            tf.clear()
            val = str(df.iloc[row_idx][col_name])
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = val
            p.font.size = Pt(data_font_size)
            p.alignment = PP_ALIGN.CENTER


# ======================
# MAIN EXECUTION
# ======================
def main():
    start_time = time.time()

    validate_files(files)
    data = load_data(files)
    plots = create_plots(data)

    print(f"💾 Files saved:")

    if OUTPUT_TYPE in ['images', 'both']:
        for key, fig in plots.items():
            new_name = TIFF_NAME_MAP[key]
            save_figure(fig, new_name)

    if OUTPUT_TYPE in ['powerpoint', 'both']:
        if os.path.exists(pptx_filename):
            os.remove(pptx_filename)

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # 🎞️ Build slides in order — just like the presentation!
        create_slide_1_title(prs)
        create_slide_2_haplotype_frequencies_and_pop_size(prs, data)
        create_slide_3_heterozygosity_and_pan_stats_grid(prs, data)
        create_slide_4_parameters_table(prs)

        prs.save(pptx_filename)
        print(f"...{pptx_filename}")

    end_time = time.time()
    execution_time = end_time - start_time
    print(f"⏱️ Execution time: {execution_time:.2f} seconds")
    print(f"📁 Output files stored in: {os.getcwd()}")


if __name__ == "__main__":
    main()
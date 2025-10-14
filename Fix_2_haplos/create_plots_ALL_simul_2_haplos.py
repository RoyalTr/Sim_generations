#!/usr/bin/env python3
"""
create_plots_ALL_simul_2_haplos.py
AUTHOR: Dr. Royal Truman
VERSION: 0.9
"""
import sys
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO
import time
import os
from datetime import datetime
from pathlib import Path  # Minor cleanup: modern path handling

start_time = time.time()

# ======================
# CONFIGURATION & CONSTANTS
# ======================

CONFIG_FILE = "config_2_haplos_v1.yaml"
DATA_FILE = 'out_2_haplos_per_gen.txt'
TITLE_IMAGE_PATH = "ALL_PER_GEN_2_HAPLOS_IMAGE.PNG"
PPTX_OUTPUT_FILENAME = 'ALL_per_gen_2_haplos.pptx'

# Column type specifications
DTYPE_SPEC = {
    'SimNr': 'int', 'Rep': 'int', 'attempt': 'int',
    'Ni': 'int', 'K': 'int', 'attempts': 'int',
    'generation': 'int', 'N': 'int'
}

# Parameter columns (used for extraction)
PARAM_COLS = ['Ni', 'r', 'K', 's_A', 'h_A', 'p_A_i', 's_B', 'h_B', 'p_B_i', 'attempts']

# Plotting configuration
PLOT_DPI = 600
PLOT_THEME = {
    'font_size_axis': 14,
    'font_size_tick': 12,
    'font_size_legend': 10,
    'grid_alpha': 0.3,
    'line_width_main': 0.8,
    'line_width_bold': 1.6,
}

# Color palette for plots
COLORS = {
    'freq_A': 'darkblue',
    'freq_Aa': 'lightblue',
    'freq_a': 'red',
    'freq_B': 'darkblue',
    'freq_Bb': 'lightblue',
    'freq_b': 'red',
    'pan_heteroz': 'darkgreen',
    'pan_homoz': 'orange',
    'population': 'black',
}

# Slide dimensions (16:9)
SLIDE_WIDTH_INCHES = 13.33
SLIDE_HEIGHT_INCHES = 7.5

# === Load YAML config file ===
try:
    import yaml
except ImportError:
    yaml = None
    print("❌ Required package 'pyyaml' not found.")
    print("👉 Install it with: pip install pyyaml")
    sys.exit(1)

# Check if config file exists
if not Path(CONFIG_FILE).exists():
    print(f"❌ Configuration file '{CONFIG_FILE}' not found in your directory.")
    print(f" Copy it over, and modify the parameters if wished, then rerun.")
    yaml = None  # Avoid unnecessary warning from PyCharm
    sys.exit(1)

# Load config
try:
    with open(CONFIG_FILE, 'r') as f:
        config = yaml.safe_load(f)
        print(f"📄 Loading configuration from '{CONFIG_FILE}'...")
except Exception as e:
    print(f"❌ Error reading or parsing '{CONFIG_FILE}': {e}")
    print("Please check that the file is valid YAML (correct indentation, colons, etc.).")
    sys.exit(1)

# Extract runtime parameters
try:
    OUTPUT_TYPE = config['OUTPUT_TYPE']
    if OUTPUT_TYPE not in ['powerpoint', 'images', 'both']:
        raise ValueError("OUTPUT_TYPE must be one of: 'powerpoint', 'images', 'both'")
except KeyError as e:
    print(f"❌ Missing required setting 'OUTPUT_TYPE' in config: {e}")
    print(f"👉 Please make sure '{CONFIG_FILE}' includes all required keys.")
    sys.exit(1)
except ValueError as e:
    print(f"❌ Invalid value in config: {e}")
    sys.exit(1)

# Load data
df = pd.read_csv(DATA_FILE, delimiter=';', dtype=DTYPE_SPEC)
print(f"📂 Reading simulation data from '{DATA_FILE}'...")
df.sort_values(['SimNr', 'Rep', 'attempt', 'generation'], inplace=True)
sim_nrs = sorted(df['SimNr'].unique())


# ========================
# Function to add title slide
# ========================
def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Top title textbox
    top_textbox = slide.shapes.add_textbox(
        left=Inches((SLIDE_WIDTH_INCHES - 9) / 2),
        top=Inches(0.5),
        width=Inches(9),
        height=Inches(0.3)
    )
    tf = top_textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "2-haplotype simulations, all trajectories. Generated with the "
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = "Sim-generations"
    run.font.size = Pt(14)
    run.font.italic = True

    run = p.add_run()
    run.text = " suite"
    run.font.size = Pt(14)

    # Bottom left
    bottom_left_textbox = slide.shapes.add_textbox(
        left=Inches(1),
        top=Inches(SLIDE_HEIGHT_INCHES) - Cm(0.5) - Inches(0.3),
        width=Inches(5),
        height=Inches(0.3)
    )
    tf = bottom_left_textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Created with create_plots_AVE_simul_2_haplos.py (v. 0.9)"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

    # Bottom right with date/time
    now = datetime.now()
    date_time_str = now.strftime("%b. %d, %Y  %H:%M")

    bottom_right_textbox = slide.shapes.add_textbox(
        left=Inches(SLIDE_WIDTH_INCHES - 1 - 5),
        top=Inches(SLIDE_HEIGHT_INCHES) - Cm(0.5) - Inches(0.3),
        width=Inches(5),
        height=Inches(0.3)
    )
    tf = bottom_right_textbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = date_time_str
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.RIGHT

    # Add image if exists
    image_path = Path(TITLE_IMAGE_PATH)
    if image_path.exists():
        top_boundary = Inches(0.5) + Inches(0.3) + Inches(1)
        bottom_boundary = Inches(SLIDE_HEIGHT_INCHES) - Cm(0.5) - Inches(0.3) - Inches(1)
        left_margin = Inches(1)
        right_margin = Inches(1)

        available_width = Inches(SLIDE_WIDTH_INCHES) - left_margin - right_margin
        available_height = bottom_boundary - top_boundary

        try:
            img_shape = slide.shapes.add_picture(
                str(image_path),
                left=left_margin,
                top=top_boundary,
                height=available_height
            )
            img_shape.left = left_margin + (available_width - img_shape.width) // 2
        except Exception as e:
            print(f"⚠️  Warning: Could not add image '{image_path}': {e}")
    else:
        print(f"⚠️  Warning: Image file '{image_path}' not found. Continuing without image.")


# ========================
# Helper function to configure common axis properties
# ========================
def configure_axis(ax, xlabel="Generations", ylabel=None, ylim=None):
    ax.set_xlabel(xlabel, fontsize=PLOT_THEME['font_size_axis'])
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=PLOT_THEME['font_size_axis'])
    if ylim:
        ax.set_ylim(ylim)
    ax.grid(True, alpha=PLOT_THEME['grid_alpha'])
    ax.tick_params(axis='both', which='major', labelsize=PLOT_THEME['font_size_tick'])
    ax.get_xaxis().set_major_formatter(plt.FuncFormatter(lambda x, p: format(int(x), ',')))


# ========================
# Plotting functions for each graph
# ========================
def plot_population_size(ax, runs_list):
    for _, run_df in runs_list:
        ax.plot(run_df['generation'], run_df['N'],
                color=COLORS['population'],
                linewidth=PLOT_THEME['line_width_main'])
    configure_axis(ax, ylabel="Population Size (N)")


def plot_freq_A(ax, runs_list, max_gen_A):
    for _, run_df in runs_list:
        truncated = run_df[run_df['generation'] <= max_gen_A]
        ax.plot(truncated['generation'], truncated['freq_A'],
                color=COLORS['freq_A'],
                linewidth=PLOT_THEME['line_width_main'],
                label="haplotype 'A'")
        ax.plot(truncated['generation'], truncated['freq_Aa'],
                color=COLORS['freq_Aa'],
                linewidth=PLOT_THEME['line_width_bold'],
                label="freq. heterozygous")
        ax.plot(truncated['generation'], truncated['freq_a'],
                color=COLORS['freq_a'],
                linewidth=PLOT_THEME['line_width_main'],
                label="haplotype 'a'")
    configure_axis(ax, ylabel="Freq. A haplotypes", ylim=(0, 1))
    handles, labels = ax.get_legend_handles_labels()
    by_label = dict(zip(labels, handles))
    ax.legend(by_label.values(), by_label.keys(),
              fontsize=PLOT_THEME['font_size_legend'], loc='best')


def plot_freq_B(ax, runs_list, max_gen_B):
    for _, run_df in runs_list:
        truncated = run_df[run_df['generation'] <= max_gen_B]
        ax.plot(truncated['generation'], truncated['freq_B'],
                color=COLORS['freq_B'],
                linewidth=PLOT_THEME['line_width_main'],
                label="haplotype 'B'")
        ax.plot(truncated['generation'], truncated['freq_Bb'],
                color=COLORS['freq_Bb'],
                linewidth=PLOT_THEME['line_width_bold'],
                label="freq. heterozygous")
        ax.plot(truncated['generation'], truncated['freq_b'],
                color=COLORS['freq_b'],
                linewidth=PLOT_THEME['line_width_main'],
                label="haplotype 'b'")
    configure_axis(ax, ylabel="Freq. B haplotypes", ylim=(0, 1))
    handles, labels = ax.get_legend_handles_labels()
    by_label = dict(zip(labels, handles))
    ax.legend(by_label.values(), by_label.keys(),
              fontsize=PLOT_THEME['font_size_legend'], loc='best')


def plot_pan_heteroz_homoz(ax, runs_list):
    for _, run_df in runs_list:
        ax.plot(run_df['generation'], run_df['pan_heteroz'],
                color=COLORS['pan_heteroz'],
                linewidth=PLOT_THEME['line_width_main'],
                label="pan heterozygous")
        ax.plot(run_df['generation'], run_df['pan_homoz'],
                color=COLORS['pan_homoz'],
                linewidth=PLOT_THEME['line_width_main'],
                label="pan homozygous")
    configure_axis(ax, ylabel="Pan heteroz and homozyg", ylim=(0, 1))
    handles, labels = ax.get_legend_handles_labels()
    by_label = dict(zip(labels, handles))
    ax.legend(by_label.values(), by_label.keys(),
              fontsize=PLOT_THEME['font_size_legend'], loc='best')


# ========================
# Function to create figure with 4 subplots (for PowerPoint)
# ========================
def create_figure(runs_list, max_gen_A, max_gen_B):
    fig, axes = plt.subplots(2, 2, figsize=(12, 9), dpi=PLOT_DPI)
    ax1, ax2, ax3, ax4 = axes.flat

    plot_population_size(ax1, runs_list)
    plot_freq_A(ax2, runs_list, max_gen_A)
    plot_freq_B(ax3, runs_list, max_gen_B)
    plot_pan_heteroz_homoz(ax4, runs_list)

    plt.tight_layout()
    return fig


# ========================
# Function to extract parameters
# ========================
def extract_parameters(first_row):
    return {
        'Ni': int(first_row['Ni']),
        'r': first_row['r'],
        'K': int(first_row['K']),
        's_A': first_row['s_A'],
        'h_A': first_row['h_A'],
        'p_A_i': first_row['p_A_i'],
        's_B': first_row['s_B'],
        'h_B': first_row['h_B'],
        'p_B_i': first_row['p_B_i'],
        'attempts': int(first_row['attempts'])
    }


# ========================
# Function to calculate max_gen for A or B
# ========================
def calculate_max_gen(group_df, runs_list, freq_col):
    max_gen = max(group_df['generation'])
    gens = []
    for _, run_df in runs_list:
        zero_gens = run_df[run_df[freq_col] == 0]['generation']
        if len(zero_gens) > 0:
            gens.append(zero_gens.iloc[0])
    if gens:
        max_gen = max(gens)
    return max_gen


# ========================
# Function to add simulation slide
# ========================
def add_simulation_slide(prs, sim_nr, params, Rep_max, total_runs, runs_list, max_gen_A, max_gen_B):
    fig = create_figure(runs_list, max_gen_A, max_gen_B)
    img_stream = BytesIO()
    fig.savefig(img_stream, format='png', dpi=PLOT_DPI, bbox_inches='tight')
    img_stream.seek(0)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = f"Simulation {sim_nr} (Repetitions = {Rep_max}, attempts = {params['attempts']}, total runs = {total_runs})"
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

    title.left = Inches(3.0)
    title.top = Inches(0.1)
    title.width = Inches(8.0)
    title.height = Inches(0.2)

    # Add parameter text box below title
    left = Inches(2)
    top = Inches(0.3)
    width = Inches(8.2)
    height = Inches(0.7)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = f"Initial population (Ni) = {params['Ni']}, growth rate (r) = {params['r']:.3f}, maximum population size (K) = {params['K']},"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

    p = tf.add_paragraph()
    p.text = f"selectivity coeff. haplotype 'A', s(A) = {params['s_A']:.3f}, dominance coeff., h(A) = {params['h_A']:.3f}, initial proportion haplotype 'A' = {params['p_A_i']:.3f},"
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

    p = tf.add_paragraph()
    p.text = f"selectivity coeff. haplotype 'B', s(B) = {params['s_B']:.3f}, dominance coeff., h(B) = {params['h_B']:.3f}, initial proportion haplotype 'B' = {params['p_B_i']:.3f}."
    p.font.size = Pt(12)
    p.alignment = PP_ALIGN.LEFT

    # Position image below text box
    img_top = top + height + Inches(0.2)
    img_height = Inches(5.8)
    slide.shapes.add_picture(img_stream, left, img_top, width=width, height=img_height)
    plt.close(fig)


# ========================
# Helper: Save individual plot to TIFF with consistent styling
# ========================
def save_plot_to_tiff(plot_func, runs_list, params, sim_nr, filename_base, max_gen=None, ylabel=None):
    """Save a single plot to TIFF with parameter overlay and consistent styling."""
    fig, ax = plt.subplots(figsize=(8, 6), dpi=PLOT_DPI)

    # Call the plotting function with or without max_gen
    if max_gen is not None:
        plot_func(ax, runs_list, max_gen)
    else:
        plot_func(ax, runs_list)

    # Ensure ylabel is set if provided (some plots set it internally, but we allow override)
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=PLOT_THEME['font_size_axis'])

    # Add parameter text box
    text = (
        f"Ni={params['Ni']}, r={params['r']:.3f}, K={params['K']}\n"
        f"s_A={params['s_A']:.3f}, h_A={params['h_A']:.3f}, p_A_i={params['p_A_i']:.3f}\n"
        f"s_B={params['s_B']:.3f}, h_B={params['h_B']:.3f}, p_B_i={params['p_B_i']:.3f}"
    )
    ax.text(0.5, 0.02, text, transform=ax.transAxes, fontsize=12,
            verticalalignment='bottom', horizontalalignment='left',
            bbox=dict(boxstyle="round,pad=0.4", facecolor="wheat", alpha=0.7))

    plt.tight_layout()

    filename = f"{filename_base}_{sim_nr}.tiff"
    fig.savefig(filename, format='tiff', dpi=PLOT_DPI, pil_kwargs={"compression": "tiff_lzw"})
    plt.close(fig)
    print(f"...{filename}")


# ========================
# Main Execution Flow
# ========================

if OUTPUT_TYPE in ['powerpoint', 'both']:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    add_title_slide(prs)

# Process each simulation
for sim_nr in sim_nrs:
    group_df = df[df['SimNr'] == sim_nr].copy()
    first_row = group_df.iloc[0]

    params = extract_parameters(first_row)
    Rep_max = group_df['Rep'].max()
    total_runs = Rep_max * params['attempts']
    runs_list = list(group_df.groupby(['Rep', 'attempt']))

    max_gen_A = calculate_max_gen(group_df, runs_list, 'freq_Aa')
    max_gen_B = calculate_max_gen(group_df, runs_list, 'freq_Bb')

    # Create PowerPoint slide
    if OUTPUT_TYPE in ['powerpoint', 'both']:
        add_simulation_slide(prs, sim_nr, params, Rep_max, total_runs, runs_list, max_gen_A, max_gen_B)

    # Save individual TIFF images
    if OUTPUT_TYPE in ['images', 'both']:
        print(f"💾 Saving files for Simulation {sim_nr}:")

        save_plot_to_tiff(
            plot_population_size,
            runs_list,
            params,
            sim_nr,
            "ALL_A_haplo_popul_size",
            ylabel="Population Size (N)"
        )

        save_plot_to_tiff(
            plot_freq_A,
            runs_list,
            params,
            sim_nr,
            "ALL_A_haplo_freq_sim",
            max_gen=max_gen_A,
            ylabel="Freq. A haplotypes"
        )

        save_plot_to_tiff(
            plot_freq_B,
            runs_list,
            params,
            sim_nr,
            "ALL_B_haplo_freq_sim",
            max_gen=max_gen_B,
            ylabel="Freq. B haplotypes"
        )

        save_plot_to_tiff(
            plot_pan_heteroz_homoz,
            runs_list,
            params,
            sim_nr,
            "ALL_Pan-het_hom_freq_sim",
            ylabel="Pan heteroz and homozyg"
        )

# Save PowerPoint
if OUTPUT_TYPE in ['powerpoint', 'both']:
    prs.save(PPTX_OUTPUT_FILENAME)
    print(f"...{PPTX_OUTPUT_FILENAME}")

# Final Output Info
current_dir = Path.cwd()
print(f"📁 Output files stored in: {current_dir}")

end_time = time.time()
execution_time = end_time - start_time
print(f"⏱️ Total runtime: {execution_time:.2f} seconds")